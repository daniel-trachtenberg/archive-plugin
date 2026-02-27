import services.filesystem_service as filesystem
import services.llm_service as llm
import services.chroma_service as chroma
import services.move_log_service as move_logs
from PyPDF2 import PdfReader
import base64
import io
import os
import logging
import json
from pptx import Presentation
import shutil
from config import settings
from datetime import datetime
from time import monotonic
import docx
import pandas as pd


def extract_text_from_pdf(file_content):
    try:
        reader = PdfReader(io.BytesIO(file_content))
        full_text = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                full_text.append(page_text.replace("\n", " "))
        return " ".join(full_text)
    except Exception as e:
        logging.error(f"Failed to extract text from PDF: {str(e)}")
        return ""


def extract_text_from_pptx(file_content):
    """
    Extract text from PowerPoint .pptx file.
    """
    try:
        prs = Presentation(io.BytesIO(file_content))
        full_text = []

        logging.info(
            f"Starting extraction from PowerPoint with {len(prs.slides)} slides"
        )

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"Slide {slide_num}:"]

            # Extract text from all shapes including titles and content placeholders
            shape_count = 0
            text_shape_count = 0
            table_count = 0
            group_shape_count = 0

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
                    text_shape_count += 1

                # Extract text from tables
                if shape.has_table:
                    table_count += 1
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text:
                                row_text.append(cell.text)
                        if row_text:
                            slide_text.append(" | ".join(row_text))

                # Extract text from group shapes
                if shape.shape_type == 6:  # GROUP shape type
                    group_shape_count += 1
                    for subshape in shape.shapes:
                        if hasattr(subshape, "text") and subshape.text:
                            slide_text.append(subshape.text)

                shape_count += 1

            # Add slide text to full text
            if (
                len(slide_text) > 1
            ):  # Only add if there's more than just the slide number
                full_text.append("\n".join(slide_text))

            logging.debug(
                f"Slide {slide_num}: Found {shape_count} shapes, {text_shape_count} with text, {table_count} tables, {group_shape_count} group shapes"
            )

        result = "\n\n".join(full_text)
        logging.info(
            f"PowerPoint extraction complete: {len(result)} characters extracted"
        )

        # Log a preview of the extracted text (first 200 chars)
        if result:
            logging.debug(f"Content preview: {result[:200]}...")

        return result
    except Exception as e:
        logging.error(f"Failed to extract text from PowerPoint: {str(e)}")
        import traceback

        logging.error(traceback.format_exc())
        return ""


def extract_text_from_docx(file_content):
    """
    Extract text from Word .docx file.
    """
    try:
        doc = docx.Document(io.BytesIO(file_content))
        full_text = []

        # Extract text from paragraphs
        for para in doc.paragraphs:
            if para.text:
                full_text.append(para.text)

        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text:
                        row_text.append(cell.text)
                if row_text:
                    full_text.append(" | ".join(row_text))

        result = "\n\n".join(full_text)
        logging.info(
            f"Word document extraction complete: {len(result)} characters extracted"
        )

        # Log a preview of the extracted text (first 200 chars)
        if result:
            logging.debug(f"Content preview: {result[:200]}...")

        return result
    except Exception as e:
        logging.error(f"Failed to extract text from Word document: {str(e)}")
        import traceback

        logging.error(traceback.format_exc())
        return ""


def extract_text_from_excel(file_content):
    """
    Extract text from Excel .xls/.xlsx files.
    """
    try:
        excel_file = io.BytesIO(file_content)

        # Try different engines if needed
        engines_to_try = ["openpyxl", "xlrd"]
        workbook = None
        sheet_names = []
        error_messages = []

        for engine in engines_to_try:
            try:
                # Try to read with specific engine
                logging.info(f"Trying to read Excel file with engine: {engine}")
                workbook = pd.ExcelFile(excel_file, engine=engine)
                sheet_names = workbook.sheet_names
                logging.info(
                    f"Successfully read Excel file with {engine}, found {len(sheet_names)} sheets"
                )
                break
            except Exception as e:
                error_messages.append(f"Engine {engine} failed: {str(e)}")
                excel_file.seek(0)  # Reset file pointer for next attempt

        full_text = []

        # If we couldn't open with pandas, try a last resort approach with raw file analysis
        if not workbook:
            error_details = "\n".join(error_messages)
            logging.error(
                f"Could not read Excel file with any engine. Errors:\n{error_details}"
            )

            # Try to get basic information about the file
            try:
                # Check if it's actually an Excel file by checking file signatures
                excel_file.seek(0)
                header = excel_file.read(8).hex()

                # Check for Excel file signatures
                if header.startswith("504b34"):  # PKZip signature (xlsx)
                    file_type = "Excel XLSX (Office Open XML)"
                elif header.startswith("d0cf11e0"):  # Compound File Binary Format (xls)
                    file_type = "Excel XLS (Binary)"
                else:
                    file_type = "Unknown (not a standard Excel format)"

                # Get file size
                excel_file.seek(0, 2)  # Seek to end
                file_size = excel_file.tell()

                full_text = [
                    "Excel file detected but could not be fully read.",
                    f"File type: {file_type}",
                    f"File size: {file_size/1024:.1f} KB",
                    "This file may be password-protected, corrupted, or use an unsupported Excel format.",
                ]

                result = "\n".join(full_text)
                return result
            except Exception as ex:
                return f"Could not read Excel file. Please ensure it's a valid Excel document. Error: {str(ex)}"

        # Add workbook metadata if available
        try:
            # Different Excel engines store properties differently
            if hasattr(workbook, "book") and hasattr(workbook.book, "properties"):
                props = workbook.book.properties
                if hasattr(props, "title") and props.title:
                    full_text.append(f"Title: {props.title}")
                if hasattr(props, "subject") and props.subject:
                    full_text.append(f"Subject: {props.subject}")
                if hasattr(props, "author") and props.author:
                    full_text.append(f"Author: {props.author}")
        except Exception as e:
            logging.warning(f"Could not extract workbook metadata: {str(e)}")

        # Add sheet names overview
        if sheet_names:
            full_text.append(
                f"Workbook contains {len(sheet_names)} sheets: {', '.join(sheet_names)}"
            )

        # Process each sheet
        sheets_processed = 0
        sheets_failed = 0

        for sheet_name in sheet_names:
            try:
                sheet_text = [f"Sheet: {sheet_name}"]

                # Read the sheet into a dataframe with the same engine
                excel_file.seek(0)
                df = pd.read_excel(
                    excel_file, sheet_name=sheet_name, engine=workbook.engine
                )

                # Skip completely empty sheets
                if df.empty:
                    sheet_text.append("(Empty sheet)")
                    full_text.append("\n".join(sheet_text))
                    sheets_processed += 1
                    continue

                # Get basic info about the sheet
                rows, cols = df.shape
                sheet_text.append(f"Dimensions: {rows} rows × {cols} columns")

                # Get column names
                if not df.columns.empty:
                    sheet_text.append(f"Columns: {', '.join(df.columns.astype(str))}")

                # Sample data - first 10 rows
                max_rows = min(10, rows)
                if max_rows > 0:
                    sheet_text.append("Data sample:")
                    rows_added = 0
                    for i in range(max_rows):
                        try:
                            row_values = df.iloc[i].astype(str)
                            # Filter out rows that just contain NaN values
                            if not all(val == "nan" for val in row_values):
                                row_text = " | ".join(row_values)
                                # Limit row text length
                                if len(row_text) > 500:
                                    row_text = row_text[:500] + "..."
                                sheet_text.append(row_text)
                                rows_added += 1
                        except Exception as row_e:
                            sheet_text.append(f"(Error reading row {i}: {str(row_e)})")

                    if rows_added == 0:
                        sheet_text.append("(No meaningful data rows found)")

                # Add sheet content to full text
                full_text.append("\n".join(sheet_text))
                sheets_processed += 1
            except Exception as sheet_e:
                logging.warning(
                    f"Error processing sheet '{sheet_name}': {str(sheet_e)}"
                )
                full_text.append(
                    f"Sheet: {sheet_name}\n(Error reading sheet: {str(sheet_e)})"
                )
                sheets_failed += 1

        # Add summary of processing
        full_text.append(
            f"\nProcessing summary: {sheets_processed} sheets processed successfully, {sheets_failed} sheets failed."
        )

        result = "\n\n".join(full_text)
        logging.info(
            f"Excel file extraction complete: {len(result)} characters extracted"
        )

        # Log a preview of the extracted text
        if result:
            logging.debug(f"Content preview: {result[:200]}...")
        else:
            logging.warning("No content extracted from Excel file")
            result = "Empty Excel workbook"

        return result
    except Exception as e:
        logging.error(f"Failed to extract text from Excel file: {str(e)}")
        import traceback

        logging.error(traceback.format_exc())
        return "Error extracting Excel content. Please check the file format."


def limit_text_for_llm(text, max_chars=8192):
    """
    Limit text size to prevent exceeding LLM context window.
    Takes the first max_chars characters from the text.

    Args:
        text (str): The input text to limit
        max_chars (int): Maximum number of characters to keep, defaults to 8192

    Returns:
        str: The limited text
    """
    if text is None:
        return ""

    if not isinstance(text, str):
        text = str(text)

    if len(text) <= max_chars:
        return text

    return text[:max_chars]


_DIRECTORY_CONTEXT_CACHE = {"value": "", "created_at": 0.0}
_DIRECTORY_CONTEXT_CACHE_TTL_SECONDS = 10
_MAX_FOLDER_PATH_DEPTH = 7


def _normalize_path_for_prompt(path: str) -> str:
    return (path or "").replace("\\", "/").strip()


def _is_visible_prompt_path(path: str) -> bool:
    normalized = _normalize_path_for_prompt(path)
    if not normalized:
        return False

    parts = [part for part in normalized.split("/") if part]
    if not parts:
        return False
    if any(part == ".chromadb" for part in parts):
        return False
    if any(part.startswith(".") for part in parts):
        return False
    return True


def _sorted_visible_paths(paths):
    visible = {
        _normalize_path_for_prompt(path)
        for path in paths
        if isinstance(path, str) and _is_visible_prompt_path(path)
    }
    return sorted(visible)


def _trim_paths_for_prompt(paths, max_items: int):
    if max_items <= 0:
        return []

    if len(paths) <= max_items:
        return list(paths)

    remaining = len(paths) - max_items
    return list(paths[:max_items]) + [
        f"... ({remaining} additional paths omitted to fit prompt limits)"
    ]


def _extract_directory_prefixes_for_prompt(paths, max_depth: int = _MAX_FOLDER_PATH_DEPTH):
    directories = set()
    for path in paths:
        normalized = _normalize_path_for_prompt(path)
        if not normalized or normalized.startswith("... ("):
            continue

        parent = os.path.dirname(normalized).replace("\\", "/").strip("/")
        if not parent:
            continue

        parts = [part for part in parent.split("/") if part]
        for idx in range(1, min(len(parts), max_depth) + 1):
            directories.add("/".join(parts[:idx]))

    return sorted(directories)


def _build_directory_context_payload(max_chars: int = 18000) -> str:
    archive_files = _sorted_visible_paths(filesystem.list_archive_files())
    indexed_files = _sorted_visible_paths(chroma.list_indexed_paths())
    existing_directories = _extract_directory_prefixes_for_prompt(archive_files)

    archive_set = set(archive_files)
    indexed_set = set(indexed_files)
    unindexed_files = sorted(archive_set - indexed_set)
    db_only_records = sorted(indexed_set - archive_set)

    tree_text = filesystem.directory_tree_for_llm(
        max_entries=max(2000, len(archive_files) + 200)
    )

    # Try increasingly smaller variants until we fit safely within max_chars.
    variants = [
        {"tree_chars": 12000, "archive": 1200, "indexed": 600, "unindexed": 1200, "db_only": 600},
        {"tree_chars": 9000, "archive": 800, "indexed": 400, "unindexed": 800, "db_only": 400},
        {"tree_chars": 7000, "archive": 500, "indexed": 250, "unindexed": 500, "db_only": 250},
        {"tree_chars": 5000, "archive": 300, "indexed": 150, "unindexed": 300, "db_only": 150},
    ]

    for variant in variants:
        payload = {
            "archive_tree": limit_text_for_llm(tree_text, max_chars=variant["tree_chars"]),
            "stats": {
                "archive_file_count": len(archive_files),
                "indexed_file_count": len(indexed_files),
                "unindexed_file_count": len(unindexed_files),
                "db_only_record_count": len(db_only_records),
            },
            "existing_directories": _trim_paths_for_prompt(
                existing_directories, variant["archive"]
            ),
            # Files currently present in archive (whether indexed or not).
            "archive_files": _trim_paths_for_prompt(archive_files, variant["archive"]),
            # Files known to Chroma index.
            "indexed_files": _trim_paths_for_prompt(indexed_files, variant["indexed"]),
            # Files present in archive but not yet in Chroma.
            "unindexed_archive_files": _trim_paths_for_prompt(
                unindexed_files, variant["unindexed"]
            ),
            # Records in Chroma that no longer exist on disk.
            "db_only_index_records": _trim_paths_for_prompt(
                db_only_records, variant["db_only"]
            ),
        }

        serialized = json.dumps(payload, separators=(",", ":"), ensure_ascii=False)
        if len(serialized) <= max_chars:
            return serialized

    minimal_payload = {
        "archive_tree": limit_text_for_llm(tree_text, max_chars=3500),
        "stats": {
            "archive_file_count": len(archive_files),
            "indexed_file_count": len(indexed_files),
            "unindexed_file_count": len(unindexed_files),
            "db_only_record_count": len(db_only_records),
        },
        "existing_directories": _trim_paths_for_prompt(existing_directories, 200),
        "unindexed_archive_files": _trim_paths_for_prompt(unindexed_files, 120),
        "db_only_index_records": _trim_paths_for_prompt(db_only_records, 120),
    }
    return json.dumps(minimal_payload, separators=(",", ":"), ensure_ascii=False)


def _invalidate_directory_context_cache():
    _DIRECTORY_CONTEXT_CACHE["value"] = ""
    _DIRECTORY_CONTEXT_CACHE["created_at"] = 0.0


def directory_structure_for_llm(force_refresh: bool = False):
    """
    Build an LLM-ready placement context with:
    - archive directory tree
    - all visible archive files
    - indexed/unindexed split vs ChromaDB
    - stale DB records
    """
    now = monotonic()
    cache_age = now - _DIRECTORY_CONTEXT_CACHE["created_at"]
    if (
        not force_refresh
        and _DIRECTORY_CONTEXT_CACHE["value"]
        and cache_age < _DIRECTORY_CONTEXT_CACHE_TTL_SECONDS
    ):
        return _DIRECTORY_CONTEXT_CACHE["value"]

    try:
        context = _build_directory_context_payload(max_chars=18000)
        _DIRECTORY_CONTEXT_CACHE["value"] = context
        _DIRECTORY_CONTEXT_CACHE["created_at"] = now
        return context
    except Exception as e:
        logging.error(f"Failed building llm directory context: {e}")
        fallback = filesystem.directory_tree_for_llm(max_entries=800)
        _DIRECTORY_CONTEXT_CACHE["value"] = fallback
        _DIRECTORY_CONTEXT_CACHE["created_at"] = now
        return fallback


def _ensure_unique_relative_path(relative_path: str) -> str:
    """
    Ensure files are never overwritten in the archive by appending numeric suffixes.
    """
    normalized = os.path.normpath(relative_path)
    full_path = os.path.join(settings.ARCHIVE_DIR, normalized)
    if not os.path.exists(full_path):
        return normalized

    parent = os.path.dirname(normalized)
    stem, ext = os.path.splitext(os.path.basename(normalized))

    counter = 1
    while True:
        candidate_name = f"{stem}_{counter}{ext}"
        candidate = os.path.join(parent, candidate_name) if parent else candidate_name
        candidate_full = os.path.join(settings.ARCHIVE_DIR, candidate)
        if not os.path.exists(candidate_full):
            return candidate
        counter += 1


def _move_trigger_for_source(source_path: str) -> str:
    if source_path.startswith("manual-upload:"):
        return "manual_upload"
    if source_path:
        return "input_watcher"
    return "plugin"


async def process_document(
    filename: str,
    content: bytes,
    source_path: str = "",
):
    try:
        logging.info(f"Processing document: {filename}")

        directory_structure = directory_structure_for_llm()

        if filename.lower().endswith(".pdf"):
            file_content = extract_text_from_pdf(content)
        elif filename.lower().endswith(".pptx"):
            file_content = extract_text_from_pptx(content)
            logging.info(
                f"Extracted PowerPoint content length: {len(file_content)} characters"
            )

            # Special handling for PowerPoint files with little or no extractable text
            if not file_content or len(file_content) < 50:
                logging.warning(
                    f"PowerPoint file {filename} has little or no extractable text"
                )
                # Use filename as a fallback for content
                basename = os.path.splitext(os.path.basename(filename))[0]
                processed_name = basename.replace("_", " ").replace("-", " ")
                file_content = f"PowerPoint presentation titled: {processed_name}"
        elif filename.lower().endswith((".docx", ".doc")):
            file_content = extract_text_from_docx(content)
            logging.info(
                f"Extracted Word document content length: {len(file_content)} characters"
            )

            # Special handling for Word files with little or no extractable text
            if not file_content or len(file_content) < 50:
                logging.warning(
                    f"Word document {filename} has little or no extractable text"
                )
                # Use filename as a fallback for content
                basename = os.path.splitext(os.path.basename(filename))[0]
                processed_name = basename.replace("_", " ").replace("-", " ")
                file_content = f"Word document titled: {processed_name}"
        elif filename.lower().endswith((".xlsx", ".xls")):
            file_content = extract_text_from_excel(content)
            logging.info(
                f"Extracted Excel file content length: {len(file_content)} characters"
            )

            # Special handling for Excel files with little or no extractable text
            if not file_content or len(file_content) < 50:
                logging.warning(
                    f"Excel file {filename} has little or no extractable text"
                )
                # Use filename as a fallback for content
                basename = os.path.splitext(os.path.basename(filename))[0]
                processed_name = basename.replace("_", " ").replace("-", " ")
                file_content = f"Excel spreadsheet titled: {processed_name}"
        else:
            file_content = content.decode("utf-8", errors="ignore")

        # Limit text size to prevent exceeding LLM context window
        file_content_for_llm = limit_text_for_llm(file_content)

        # Step 1: Generate semantic summary used for organization and retrieval.
        logging.info(f"Generating summary for: {filename}")
        file_summary = await llm.get_file_summary(
            filename=filename,
            content=file_content_for_llm,
        )
        logging.info(f"Generated summary: {file_summary}")

        # Step 2: Get path suggestion based on the summary and directory structure
        logging.info(f"Getting path suggestion for: {filename} based on summary")
        suggested_path = await llm.get_path_from_summary(
            filename=filename,
            summary=file_summary,
            directory_structure=directory_structure,
        )
        logging.info(f"Initial suggested path: {suggested_path}")

        # Sanitize the suggested path to ensure proper file placement
        suggested_path = sanitize_path_suggestion(suggested_path, filename)

        # Build final file path and avoid accidental overwrites.
        suggested_path = os.path.join(suggested_path, filename)

        # Clean up duplicate path segments
        path_parts = suggested_path.replace("\\", "/").split("/")
        corrected_path_parts = [
            path_parts[i]
            for i in range(len(path_parts))
            if i == 0 or path_parts[i] != path_parts[i - 1]
        ]

        final_path = os.path.normpath("/".join(corrected_path_parts))
        final_path = _ensure_unique_relative_path(final_path)
        logging.info(f"Final document path: {final_path}")

        # Save file to filesystem
        if not filesystem.save_file(content, final_path):
            raise RuntimeError("Failed to save document to archive filesystem.")
        _invalidate_directory_context_cache()

        # Add to vector database
        embedding_payload = (
            f"Filename: {filename}\n"
            f"Summary: {file_summary or 'N/A'}\n"
            f"Content: {file_content_for_llm}"
        )
        chroma.add_document_to_collection(final_path, embedding_payload)

        # Log the final path to the terminal
        print(f"Document moved to: {final_path}")

        destination_path = os.path.join(settings.ARCHIVE_DIR, final_path)
        move_logs.record_move(
            source_path=source_path or filename,
            destination_path=destination_path,
            item_type="file",
            trigger=_move_trigger_for_source(source_path),
            status="success",
        )

        logging.info(f"Successfully processed: {filename}")
        return final_path
    except Exception as e:
        logging.error(f"Error processing: {filename}. Error: {str(e)}")
        move_logs.record_move(
            source_path=source_path or filename,
            destination_path="",
            item_type="file",
            trigger=_move_trigger_for_source(source_path),
            status="failed",
            note=str(e)[:500],
        )
        return None


async def process_image(
    filename: str,
    content: bytes,
    source_path: str = "",
):
    try:
        logging.info(f"Processing image: {filename}")

        directory_structure = directory_structure_for_llm()

        # No need to encode the image for CLIP analysis - we'll use binary content directly
        # The encoded version is still needed for some other potential uses
        encoded_image = base64.b64encode(content).decode("utf-8")

        file_extension = filename.split(".")[-1].lower()
        if file_extension in ["jpg", "jpeg"]:
            media_type = "image/jpeg"
        elif file_extension == "png":
            media_type = "image/png"
        elif file_extension == "gif":
            media_type = "image/gif"
        else:
            media_type = "application/octet-stream"

        # Try to get CLIP description directly
        image_summary = ""
        suggested_path = ""
        try:
            # Get the image summary from the CLIP model and LLM
            image_summary = await llm.get_image_summary(
                filename=filename,
                encoded_image=encoded_image,
                media_type=media_type,
            )
            logging.info(f"Image summary: {image_summary}")
        except Exception as e:
            logging.error(
                f"Error getting image summary, falling back to direct path suggestion: {e}"
            )
            # Use the path suggestion for image directly on failure
            suggested_path = await llm.get_path_suggestion_for_image(
                filename=filename,
                encoded_image=encoded_image,
                directory_structure=directory_structure,
                media_type=media_type,
            )
            logging.info(f"Got path suggestion directly: {suggested_path}")

        # If we got a summary, use it to get a path suggestion
        if image_summary:
            # Get path suggestion based on the image summary
            logging.info(
                f"Getting path suggestion for image: {filename} based on summary"
            )
            suggested_path = await llm.get_path_from_summary(
                filename=filename,
                summary=image_summary,
                directory_structure=directory_structure,
            )
            logging.info(f"Path suggestion from summary: {suggested_path}")

        # Sanitize the suggested path to ensure proper file placement
        suggested_path = sanitize_path_suggestion(suggested_path, filename)

        # Build final file path and avoid accidental overwrites.
        suggested_path = os.path.join(suggested_path, filename)

        # Clean up duplicate path segments
        path_parts = suggested_path.replace("\\", "/").split("/")
        corrected_path_parts = [
            path_parts[i]
            for i in range(len(path_parts))
            if i == 0 or path_parts[i] != path_parts[i - 1]
        ]

        final_path = os.path.normpath("/".join(corrected_path_parts))
        final_path = _ensure_unique_relative_path(final_path)
        logging.info(f"Final image path: {final_path}")

        # Save file to filesystem
        if not filesystem.save_file(content, final_path):
            raise RuntimeError("Failed to save image to archive filesystem.")
        _invalidate_directory_context_cache()

        # Add to vector database
        chroma.add_image_to_collection(
            final_path,
            content,
            summary=f"Filename: {filename}\nSummary: {image_summary or filename}",
        )

        # Log the final path to the terminal
        print(f"Image moved to: {final_path}")

        destination_path = os.path.join(settings.ARCHIVE_DIR, final_path)
        move_logs.record_move(
            source_path=source_path or filename,
            destination_path=destination_path,
            item_type="file",
            trigger=_move_trigger_for_source(source_path),
            status="success",
        )

        logging.info(f"Successfully processed: {filename}")
        return final_path
    except Exception as e:
        logging.error(f"Error processing: {filename}. Error: {str(e)}")
        move_logs.record_move(
            source_path=source_path or filename,
            destination_path="",
            item_type="file",
            trigger=_move_trigger_for_source(source_path),
            status="failed",
            note=str(e)[:500],
        )
        return None


async def process_folder(
    folder_name: str,
    folder_path: str,
):
    """
    Process a folder and move it to the suggested path within the Archive folder.
    Uses the folder name and filenames inside as content for classification.

    Args:
        folder_name: The name of the folder
        folder_path: The full path to the folder in the Input directory

    Returns:
        The final path where the folder was moved to in the Archive
    """
    try:
        logging.info(f"Processing folder: {folder_name} at path: {folder_path}")

        # First verify the folder still exists
        if not os.path.exists(folder_path):
            logging.error(f"Folder no longer exists at {folder_path}, cannot process")
            return None

        # Verify it's actually a directory
        if not os.path.isdir(folder_path):
            logging.error(f"Path {folder_path} is not a directory, cannot process")
            return None

        directory_structure = directory_structure_for_llm()

        # Create enhanced content with detailed folder analysis
        folder_content = f"FOLDER ANALYSIS:\n\nFolder name: {folder_name}\n\n"

        # Track file types for better categorization
        file_counts = {
            "images": 0,
            "documents": 0,
            "spreadsheets": 0,
            "presentations": 0,
            "audio": 0,
            "video": 0,
            "code": 0,
            "data": 0,
            "archives": 0,
            "other": 0,
        }

        file_extensions = []
        total_files = 0
        total_subfolders = 0

        try:
            # First gather statistics about the folder contents
            for root, dirs, files in os.walk(folder_path):
                # Count subfolders at the top level
                if root == folder_path:
                    total_subfolders = len(dirs)

                total_files += len(files)

                for file in files:
                    if file.startswith("."):
                        continue

                    # Track file extension
                    ext = os.path.splitext(file)[1].lower()
                    if ext and ext not in file_extensions:
                        file_extensions.append(ext)

                    # Count file types
                    if ext in [
                        ".jpg",
                        ".jpeg",
                        ".png",
                        ".gif",
                        ".bmp",
                        ".tiff",
                        ".webp",
                    ]:
                        file_counts["images"] += 1
                    elif ext in [".doc", ".docx", ".txt", ".rtf", ".odt", ".pdf"]:
                        file_counts["documents"] += 1
                    elif ext in [".xls", ".xlsx", ".csv", ".ods"]:
                        file_counts["spreadsheets"] += 1
                    elif ext in [".ppt", ".pptx", ".odp"]:
                        file_counts["presentations"] += 1
                    elif ext in [".mp3", ".wav", ".ogg", ".flac", ".aac", ".m4a"]:
                        file_counts["audio"] += 1
                    elif ext in [".mp4", ".mov", ".avi", ".mkv", ".webm", ".flv"]:
                        file_counts["video"] += 1
                    elif ext in [
                        ".py",
                        ".js",
                        ".html",
                        ".css",
                        ".java",
                        ".c",
                        ".cpp",
                        ".php",
                        ".rb",
                        ".go",
                        ".ts",
                    ]:
                        file_counts["code"] += 1
                    elif ext in [".json", ".xml", ".yaml", ".sql", ".db"]:
                        file_counts["data"] += 1
                    elif ext in [".zip", ".rar", ".tar", ".gz", ".7z"]:
                        file_counts["archives"] += 1
                    else:
                        file_counts["other"] += 1

            # Add file type summary
            folder_content += f"File type summary:\n"
            for file_type, count in file_counts.items():
                if count > 0:
                    folder_content += f"- {file_type}: {count} files\n"

            if file_extensions:
                folder_content += f"\nFile extensions: {', '.join(file_extensions)}\n"

            folder_content += f"\nTotal files: {total_files}\n"
            folder_content += f"Total subfolders: {total_subfolders}\n\n"

            # Now list specific files and folders for context
            folder_content += "FOLDER CONTENTS:\n\n"

            # Get list of files for classification
            file_list = []
            subfolder_list = []

            # Just get the top level contents for conciseness
            for item in os.listdir(folder_path):
                item_path = os.path.join(folder_path, item)

                # Skip hidden files/folders
                if item.startswith("."):
                    continue

                if os.path.isdir(item_path):
                    subfolder_list.append(item)
                else:
                    file_list.append(item)

            # Add folder content details
            if subfolder_list:
                folder_content += "Subfolders:\n"
                for subfolder in sorted(subfolder_list):
                    subfolder_desc = subfolder.replace("_", " ").replace("-", " ")
                    folder_content += f"- {subfolder_desc}\n"
                folder_content += "\n"

            if file_list:
                folder_content += "Files:\n"
                for file in sorted(file_list):
                    file_desc = file.replace("_", " ").replace("-", " ")
                    folder_content += f"- {file_desc}\n"
            else:
                folder_content += "- (Empty folder)\n"

        except Exception as e:
            logging.error(f"Error analyzing folder contents: {str(e)}")
            folder_content += "Error reading folder contents"

        # Limit text size to prevent exceeding LLM context window
        folder_content_for_llm = limit_text_for_llm(folder_content)

        # Step 1: Generate a summary of the folder content
        logging.info(f"Generating summary for folder: {folder_name}")
        folder_summary = await llm.get_file_summary(
            filename=folder_name,
            content=folder_content_for_llm,
        )
        logging.info(f"Generated folder summary: {folder_summary}")

        # Step 2: Get path suggestion based on the folder summary
        logging.info(
            f"Getting path suggestion for folder: {folder_name} based on summary"
        )
        suggested_path = await llm.get_path_from_summary(
            filename=folder_name,
            summary=folder_summary,
            directory_structure=directory_structure,
        )

        logging.info(f"Initial suggested path for folder: {suggested_path}")

        # Add the folder name to the path
        if suggested_path:
            final_path = os.path.join(suggested_path, folder_name)
        else:
            final_path = folder_name

        final_path = os.path.normpath(final_path)

        # Get the full source and destination paths
        source_path = folder_path
        dest_path = os.path.join(settings.ARCHIVE_DIR, final_path)

        logging.info(f"Planning to move folder from {source_path} to {dest_path}")

        # Verify again that the source folder exists before copying
        if not os.path.exists(source_path):
            logging.error(
                f"Source folder no longer exists at {source_path}, cannot copy"
            )
            return None

        # Create a temporary copy to ensure we don't lose the folder during processing
        temp_copy_path = os.path.join(
            settings.ARCHIVE_DIR,
            f"temp_{folder_name}_{int(datetime.now().timestamp())}",
        )
        try:
            # Create a safe copy of the folder first
            logging.info(f"Creating temporary copy at {temp_copy_path}")
            shutil.copytree(source_path, temp_copy_path)
        except Exception as e:
            logging.error(f"Error creating temporary copy of folder: {str(e)}")
            return None

        # If the destination exists but isn't a directory, find an alternative
        if os.path.exists(dest_path) and not os.path.isdir(dest_path):
            i = 1
            original_path = dest_path
            while os.path.exists(dest_path) and not os.path.isdir(dest_path):
                new_folder_name = f"{folder_name}_{i}"
                final_path = os.path.join(suggested_path, new_folder_name)
                dest_path = os.path.join(settings.ARCHIVE_DIR, final_path)
                i += 1
            logging.info(
                f"Destination exists as file, using alternative path: {dest_path}"
            )

        # Ensure the parent directory exists for the destination
        parent_dir = os.path.dirname(dest_path)
        os.makedirs(parent_dir, exist_ok=True)

        # Move from the temporary copy to the final destination
        try:
            if os.path.exists(dest_path):
                if os.path.isdir(dest_path):
                    # If destination exists and is a directory, merge contents
                    logging.info(f"Destination exists, merging contents")
                    for item in os.listdir(temp_copy_path):
                        src = os.path.join(temp_copy_path, item)
                        dst = os.path.join(dest_path, item)
                        if os.path.isdir(src):
                            if not os.path.exists(dst):
                                shutil.copytree(src, dst)
                            else:
                                # Recursively merge subdirectories
                                for subitem in os.listdir(src):
                                    src_sub = os.path.join(src, subitem)
                                    dst_sub = os.path.join(dst, subitem)
                                    if os.path.isdir(src_sub):
                                        if not os.path.exists(dst_sub):
                                            shutil.copytree(src_sub, dst_sub)
                                    elif not os.path.exists(dst_sub):
                                        shutil.copy2(src_sub, dst_sub)
                                    else:
                                        # Handle duplicate files by creating a unique name
                                        i = 1
                                        name, ext = os.path.splitext(subitem)
                                        while os.path.exists(dst_sub):
                                            new_name = f"{name}_{i}{ext}"
                                            dst_sub = os.path.join(dst, new_name)
                                            i += 1
                                        shutil.copy2(src_sub, dst_sub)
                        elif not os.path.exists(dst):
                            shutil.copy2(src, dst)
                        else:
                            # Handle duplicate files by creating a unique name
                            i = 1
                            name, ext = os.path.splitext(item)
                            while os.path.exists(dst):
                                new_name = f"{name}_{i}{ext}"
                                dst = os.path.join(dest_path, new_name)
                                i += 1
                            shutil.copy2(src, dst)
                else:
                    # Edge case: destination exists but is not a directory
                    logging.error(
                        f"Destination exists but is not a directory: {dest_path}"
                    )
                    # Create a new unique folder name
                    i = 1
                    while os.path.exists(dest_path):
                        new_folder_name = f"{folder_name}_{i}"
                        final_path = os.path.join(suggested_path, new_folder_name)
                        dest_path = os.path.join(settings.ARCHIVE_DIR, final_path)
                        i += 1
                    os.makedirs(dest_path, exist_ok=True)
                    # Copy from temp to new location
                    for item in os.listdir(temp_copy_path):
                        src = os.path.join(temp_copy_path, item)
                        dst = os.path.join(dest_path, item)
                        if os.path.isdir(src):
                            if not os.path.exists(dst):
                                shutil.copytree(src, dst)
                            else:
                                # Handle duplicate folder by creating a unique name
                                i = 1
                                while os.path.exists(dst):
                                    new_name = f"{item}_{i}"
                                    dst = os.path.join(dest_path, new_name)
                                    i += 1
                                shutil.copytree(src, dst)
                        else:
                            if not os.path.exists(dst):
                                shutil.copy2(src, dst)
                            else:
                                # Handle duplicate files by creating a unique name
                                i = 1
                                name, ext = os.path.splitext(item)
                                while os.path.exists(dst):
                                    new_name = f"{name}_{i}{ext}"
                                    dst = os.path.join(dest_path, new_name)
                                    i += 1
                                shutil.copy2(src, dst)
            else:
                # Destination doesn't exist, move the temp folder to destination
                logging.info(f"Creating new directory at destination")
                # Create parent directories
                os.makedirs(os.path.dirname(dest_path), exist_ok=True)
                # Move the temp folder to the final destination
                shutil.move(temp_copy_path, dest_path)
        except Exception as e:
            logging.error(f"Error moving folder to final destination: {str(e)}")
            # Try to clean up temporary files
            if os.path.exists(temp_copy_path):
                try:
                    shutil.rmtree(temp_copy_path)
                except:
                    pass
            return None

        # Clean up temporary files
        if os.path.exists(temp_copy_path) and temp_copy_path != dest_path:
            try:
                shutil.rmtree(temp_copy_path)
            except Exception as e:
                logging.warning(
                    f"Could not remove temporary folder {temp_copy_path}: {str(e)}"
                )

        # Log the final path to the terminal
        print(f"Folder moved to: {final_path}")
        _invalidate_directory_context_cache()

        # Now that the folder is fully processed and in its final location,
        # trigger a reconciliation to update the database with the new files
        print(f"Updating database with the new files...")
        try:
            # We'll use a targeted approach to only update this specific folder
            # rather than running a full reconciliation
            # Get all files in the directory that was moved
            files_to_process = []
            for root, _, files in os.walk(dest_path):
                for file in files:
                    if file.startswith("."):
                        continue
                    file_path = os.path.join(root, file)
                    rel_path = os.path.relpath(file_path, settings.ARCHIVE_DIR)
                    files_to_process.append(rel_path)

            print(f"Found {len(files_to_process)} files to add to the database")

            # Persist per-file move logs for this folder ingestion.
            move_entries = []
            for rel_path in files_to_process:
                destination_file = os.path.join(settings.ARCHIVE_DIR, rel_path)
                relative_inside_folder = os.path.relpath(destination_file, dest_path)
                source_candidate = os.path.normpath(
                    os.path.join(source_path, relative_inside_folder)
                )

                move_entries.append(
                    {
                        "source_path": (
                            source_candidate if os.path.exists(source_candidate) else source_path
                        ),
                        "destination_path": destination_file,
                        "item_type": "file",
                        "trigger": "input_watcher",
                        "status": "success",
                        "note": f"folder:{folder_name}",
                    }
                )

            if move_entries:
                move_logs.record_moves(move_entries)

            # Process each file and add it to ChromaDB
            for file_path in files_to_process:
                try:
                    content = filesystem.fetch_content(file_path)
                    if content:
                        is_image = file_path.lower().endswith(
                            (".jpg", ".jpeg", ".png", ".gif", ".webp")
                        )
                        if is_image:
                            chroma.add_image_to_collection(file_path, content)
                        else:
                            # Extract text based on file type
                            text_content = extract_text_for_file_type(
                                file_path, content
                            )
                            chroma.add_document_to_collection(file_path, text_content)
                except Exception as e:
                    logging.error(
                        f"Error adding file to database: {file_path}. Error: {str(e)}"
                    )

            print(f"✓ Database updated with new files")
        except Exception as e:
            logging.error(f"Error updating database after folder processing: {str(e)}")

        move_logs.record_move(
            source_path=source_path,
            destination_path=dest_path,
            item_type="folder",
            trigger="input_watcher",
            status="success",
            note=f"files:{len(files_to_process) if 'files_to_process' in locals() else 0}",
        )

        logging.info(f"Successfully processed folder: {folder_name}")
        return final_path
    except Exception as e:
        logging.error(f"Error processing folder: {folder_name}. Error: {str(e)}")
        move_logs.record_move(
            source_path=folder_path,
            destination_path="",
            item_type="folder",
            trigger="input_watcher",
            status="failed",
            note=str(e)[:500],
        )
        return None


def sanitize_path_suggestion(suggested_path, filename):
    """
    Sanitize path suggestions to ensure correct file placement.

    Fixes common issues like:
    1. Paths containing existing filenames (creating file-inside-file situations)
    2. File extensions in directory names
    3. Missing directory structure

    Args:
        suggested_path (str): The path suggested by the LLM
        filename (str): The name of the file being processed

    Returns:
        str: A sanitized path that ensures proper file placement
    """
    raw_path = (suggested_path or "").strip()
    logging.info(f"Sanitizing path suggestion: {raw_path} for file {filename}")

    # First, normalize path separators
    suggested_path = raw_path.replace("\\", "/")

    # Check if the suggested path already ends with the filename
    if suggested_path.endswith(filename):
        # Path already includes filename - extract the directory part
        directory_path = os.path.dirname(suggested_path)
        logging.info(
            f"Path already includes filename, extracted directory: {directory_path}"
        )
        return directory_path

    # Get the base filename without path
    basename = os.path.basename(filename)

    # Check if any path component resembles a filename (contains periods)
    path_parts = suggested_path.split("/")
    i = 0
    while i < len(path_parts):
        part = path_parts[i]

        # Skip empty parts
        if not part:
            i += 1
            continue

        # Check if this component is or contains our filename
        if part == basename or filename in part:
            # Remove this part as it's the filename we're trying to place
            path_parts.pop(i)
            logging.info(f"Removed filename from path parts at position {i}")
            continue

        # Check if this part looks like a filename with extension
        if "." in part and not part.startswith("."):
            # Check if it's a common file extension pattern
            ext = os.path.splitext(part)[1].lower()
            if ext and len(ext) <= 5:  # .html, .jpeg, .docx, etc.
                # This looks like a filename - remove it entirely
                path_parts.pop(i)
                logging.info(f"Removed file-like component: {part}")
                continue

        # Remove characters that are invalid or noisy in folder names.
        part = "".join(c for c in part if c.isalnum() or c in (" ", "_", "-"))
        part = part.strip().replace("  ", " ")
        path_parts[i] = part
        i += 1

    # Rebuild the path
    sanitized_parts = [segment for segment in path_parts if segment]
    sanitized_path = "/".join(sanitized_parts[:_MAX_FOLDER_PATH_DEPTH])

    # Make sure we don't have an empty path
    if not sanitized_path:
        file_extension = os.path.splitext(filename)[1].lower()
        if file_extension in [".jpg", ".jpeg", ".png", ".gif", ".webp", ".heic"]:
            sanitized_path = "Images"
        elif file_extension in [".pdf", ".doc", ".docx", ".txt", ".md", ".rtf"]:
            sanitized_path = "Documents"
        elif file_extension in [".csv", ".xls", ".xlsx"]:
            sanitized_path = "Data"
        elif file_extension in [".mp3", ".wav", ".flac"]:
            sanitized_path = "Music"
        elif file_extension in [".mp4", ".mov", ".avi"]:
            sanitized_path = "Videos"
        else:
            sanitized_path = "Files"

    logging.info(f"Sanitized path: {sanitized_path}")
    return sanitized_path


def extract_text_for_file_type(file_path, content):
    """Helper function to extract text based on file type"""
    try:
        if file_path.lower().endswith(".pdf"):
            return extract_text_from_pdf(content)
        elif file_path.lower().endswith(".pptx"):
            return extract_text_from_pptx(content)
        elif file_path.lower().endswith((".docx", ".doc")):
            return extract_text_from_docx(content)
        elif file_path.lower().endswith((".xlsx", ".xls")):
            return extract_text_from_excel(content)
        else:
            # Try to decode as text
            try:
                return content.decode("utf-8", errors="ignore")
            except:
                return f"Binary file: {os.path.basename(file_path)}"
    except Exception as e:
        logging.error(f"Error extracting text for {file_path}: {str(e)}")
        return f"Error extracting content from {os.path.basename(file_path)}"


async def reconcile_filesystem_with_chroma():
    """
    Reconcile the filesystem with the ChromaDB database.
    This function scans the Archive directory and ensures ChromaDB reflects its current state.
    """
    try:
        print("\n========== STARTING DATABASE RECONCILIATION ==========")
        logging.info("Starting filesystem and ChromaDB reconciliation...")

        # Get all files in the Archive directory
        def get_all_files(directory):
            all_files = []
            for root, _, files in os.walk(directory):
                # Skip the entire ChromaDB directory and its subdirectories
                if ".chromadb" in root.split(os.sep):
                    continue

                for file in files:
                    # Skip hidden files
                    if file.startswith("."):
                        continue
                    full_path = os.path.join(root, file)
                    relative_path = os.path.relpath(full_path, settings.ARCHIVE_DIR)
                    all_files.append(relative_path)
            return all_files

        # Get files from file system
        filesystem_files = set(get_all_files(settings.ARCHIVE_DIR))
        print(f"Found {len(filesystem_files)} files in filesystem")

        # Get all document IDs from ChromaDB (these are the file paths)
        try:
            collection = chroma.ensure_collection_exists()
            if not collection:
                logging.error("Could not access ChromaDB collection")
                print("ERROR: Could not access ChromaDB collection")
                return False

            # Get all document IDs from the collection
            chroma_content = collection.get()
            chroma_files = set(
                chroma_content["ids"]
                if chroma_content
                and "ids" in chroma_content
                and len(chroma_content["ids"]) > 0
                else []
            )
            print(f"Found {len(chroma_files)} files in ChromaDB")
        except Exception as e:
            logging.error(f"Error getting files from ChromaDB: {str(e)}")
            print(f"ERROR: Failed to get files from ChromaDB: {str(e)}")
            chroma_files = set()

        # Files that exist in filesystem but not in ChromaDB need to be added
        files_to_add = filesystem_files - chroma_files
        print(f"Files to add to ChromaDB: {len(files_to_add)}")

        if files_to_add:
            print("\n--- Adding missing files to ChromaDB ---")

        added_count = 0
        for file_path in files_to_add:
            try:
                # Skip any ChromaDB internal files that might have been missed
                if ".chromadb" in file_path.split(os.sep):
                    continue

                content = filesystem.fetch_content(file_path)
                if content:
                    is_image = file_path.lower().endswith(
                        (".jpg", ".jpeg", ".png", ".gif", ".webp")
                    )

                    if is_image:
                        chroma.add_image_to_collection(file_path, content)
                        print(f"✓ Added image: {file_path}")
                        logging.info(f"Added image to ChromaDB: {file_path}")
                        added_count += 1
                    else:
                        text_content = extract_text_for_file_type(file_path, content)
                        chroma.add_document_to_collection(file_path, text_content)
                        print(f"✓ Added document: {file_path}")
                        logging.info(f"Added document to ChromaDB: {file_path}")
                        added_count += 1
                else:
                    print(f"✗ Skipped file (could not read content): {file_path}")
            except Exception as e:
                logging.error(
                    f"Error adding file to ChromaDB during reconciliation: {file_path}, {str(e)}"
                )
                print(f"✗ Failed to add: {file_path} - {str(e)}")

        # Files that exist in ChromaDB but not in filesystem need to be removed
        files_to_remove = chroma_files - filesystem_files
        print(f"\nFiles to remove from ChromaDB: {len(files_to_remove)}")

        if files_to_remove:
            print("\n--- Removing obsolete files from ChromaDB ---")

        removed_count = 0
        for file_path in files_to_remove:
            try:
                # Skip any ChromaDB internal files that might have been included in the ChromaDB IDs
                if ".chromadb" in file_path.split(os.sep):
                    continue

                chroma.delete_item(file_path)
                print(f"✓ Removed: {file_path}")
                logging.info(f"Removed file from ChromaDB: {file_path}")
                removed_count += 1
            except Exception as e:
                logging.error(
                    f"Error removing file from ChromaDB during reconciliation: {file_path}, {str(e)}"
                )
                print(f"✗ Failed to remove: {file_path} - {str(e)}")

        print(f"\n========== RECONCILIATION COMPLETE ==========")
        print(f"Added: {added_count} files, Removed: {removed_count} files")
        print(f"Current database status: {len(filesystem_files)} files indexed\n")

        if added_count or removed_count:
            _invalidate_directory_context_cache()

        logging.info(
            f"Reconciliation complete. Added {added_count} files, removed {removed_count} files."
        )
        return True
    except Exception as e:
        logging.error(f"Error during reconciliation: {str(e)}")
        print(f"ERROR: Reconciliation failed: {str(e)}")
        return False
